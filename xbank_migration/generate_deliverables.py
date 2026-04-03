from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, ListFlowable, ListItem
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT

BASE = Path('/Users/baykus/.openclaw/workspace/xbank_migration')
ACCENT = RGBColor(230, 126, 34)
DARK = RGBColor(34, 34, 34)
MUTED = RGBColor(90, 90, 90)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.0), Inches(1.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = DARK
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(1, Inches(0.8), Inches(2.1), Inches(2.5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(10.5), Inches(1.2))
    stf = sub.text_frame
    p2 = stf.paragraphs[0]
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(15)
    r2.font.color.rgb = MUTED
    return slide


def add_bullets_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.2), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = DARK

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(10.7), Inches(5.3))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.text = b
        par.level = 0
        par.font.size = Pt(18)
        par.font.color.rgb = DARK
        par.space_after = Pt(12)
    if note:
        nb = slide.shapes.add_textbox(Inches(0.9), Inches(6.5), Inches(10.5), Inches(0.6))
        np = nb.text_frame.paragraphs[0]
        nr = np.add_run()
        nr.text = note
        nr.font.size = Pt(11)
        nr.font.italic = True
        nr.font.color.rgb = MUTED
    return slide


def build_internal_pptx(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs, 'X Bank Merchant Migration Fırsatı', 'Finhouse iç değerlendirme sunumu')
    add_bullets_slide(prs, 'Konu Nedir?', [
        'X Bank, Payten benzeri dış altyapıdan kendi sanal POS sistemine geçiyor.',
        'İlk 60 merchant taşınmış durumda; sırada yaklaşık 18.000 merchant var.',
        'Bu iş teknik entegrasyondan çok operasyon + migration programı.'
    ])
    add_bullets_slide(prs, 'Neden Fırsat?', [
        'Büyük ölçekli delivery işi ve güçlü banka referansı yaratır.',
        'Uzun süreli operasyon ve işlem bazlı gelir potansiyeli vardır.',
        'Migration sonrası ek ürün ve hizmet fırsatları doğabilir.'
    ])
    add_bullets_slide(prs, 'Ana Zorluklar', [
        'Merchant profilleri homojen değil: hazır paket, custom, ajans bağımlı yapılar var.',
        'Herkesin teknik ekibi yok; aktif takip ve yakın yönlendirme gerekiyor.',
        'Toplu duyuru ile geçiş olmaz; segment bazlı operasyon şart.'
    ])
    add_bullets_slide(prs, 'Bizim Avantajımız', [
        'Payment / fintech bağlamını biliyoruz.',
        'QNB benzeri migration deneyimi referans değeri taşıyor.',
        'Operasyon + teknik + ticari dili birlikte kurabiliyoruz.'
    ])
    add_bullets_slide(prs, 'Maliyet ve Ticari Okuma', [
        'Yeni çalışmada baz senaryoda merchant başı maliyet ~61 USD görünüyor.',
        '65 USD bandı mümkün ama marj sınırlı.',
        'Sadece setup fee ile gitmek riskli; hibrit model daha sağlıklı.'
    ], note='Öneri: discovery/pilot sabit bedeli + merchant başı dönüşüm bedeli + işlem bazlı gelir')
    add_bullets_slide(prs, 'Önerilen Strateji', [
        'Migration factory yaklaşımı sunalım.',
        'AI destekli delivery modeli ve hazır paket geçiş kitleri kuralım.',
        'İlk aşamada az ama doğru veri isteyelim; discovery ile derinleşelim.'
    ])
    add_bullets_slide(prs, 'Sonraki Adım', [
        'İç hizalanma ve teklif çerçevesinin netleşmesi.',
        'X Bank için sade müşteri deck’i ile ilk temas.',
        'Discovery fazı ve veri talebinin kontrollü açılması.'
    ])
    prs.save(out_path)


def build_client_pptx(out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_title_slide(prs, 'X Bank Sanal POS Merchant Migration Programı', 'Kontrollü, ölçeklenebilir ve düşük sürtünmeli geçiş yaklaşımı')
    add_bullets_slide(prs, 'Yönetici Özeti', [
        'X Bank yeni sanal POS platformunu devreye almış durumda.',
        'İlk geçişler önemli bir başlangıç sağlıyor.',
        'Sıradaki faz, daha geniş merchant portföyünün kontrollü şekilde dönüştürülmesi.'
    ])
    add_bullets_slide(prs, 'Yaklaşımımız', [
        'Merchant segmentasyonu',
        'Dalga bazlı migration planı',
        'Operasyonel takip ve iletişim yönetimi',
        'Teknik onboarding koordinasyonu'
    ])
    add_bullets_slide(prs, 'Neden Bu Model?', [
        'Merchant portföyleri teknik açıdan homojen değildir.',
        'Başarılı migration için tek tip değil, segment bazlı model gerekir.',
        'Bu yaklaşım hem hız hem de kontrol sağlar.'
    ])
    add_bullets_slide(prs, 'Önerilen Çalışma Modeli', [
        '1. Ön değerlendirme ve discovery',
        '2. Pilot ve playbook doğrulama',
        '3. Dalga bazlı yaygınlaştırma ve hypercare'
    ])
    add_bullets_slide(prs, 'İlk Aşamada Gerekli Veri', [
        'Toplam ve aktif merchant sayısı',
        'İşlem hacmi / adet özeti',
        'Entegrasyon tipi veya platform dağılımı',
        'İlk geçişlerden çıkan öğrenimler',
        'Hedef takvim ve başarı kriteri'
    ], note='İlk değerlendirme için sınırlı ama karar verdirici veriyle ilerlemek yeterlidir.')
    add_bullets_slide(prs, 'Beklenen Çıktılar', [
        'Segment bazlı migration stratejisi',
        'Pilot kapsamı',
        'Ekip ve kapasite planı',
        'Risk matrisi ve önerilen ticari model'
    ])
    add_bullets_slide(prs, 'Finhouse Katkısı', [
        'Ödeme sistemleri ve saha gerçekliğini birlikte okuyan yaklaşım',
        'Operasyon ve teknik enablement’i aynı çatı altında ele alma',
        'Ölçülebilir, aşamalı ve kontrollü delivery modeli'
    ])
    add_bullets_slide(prs, 'Sonraki Adım', [
        'Kısa discovery oturumu',
        'Ön veri setinin paylaşılması',
        'Pilot kapsamının birlikte belirlenmesi',
        'Detaylı çalışma planı ve teklif yapısının netleşmesi'
    ])
    prs.save(out_path)


def build_pdf(out_path, title, sections):
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='CustomTitle', fontSize=20, leading=24, spaceAfter=16, textColor='#222222'))
    styles.add(ParagraphStyle(name='CustomH', fontSize=14, leading=18, spaceBefore=10, spaceAfter=8, textColor='#E67E22'))
    styles.add(ParagraphStyle(name='CustomBody', fontSize=10.5, leading=15, alignment=TA_LEFT))
    doc = SimpleDocTemplate(str(out_path), pagesize=A4, leftMargin=40, rightMargin=40, topMargin=40, bottomMargin=40)
    story = [Paragraph(title, styles['CustomTitle']), Spacer(1, 8)]
    for sec_title, bullets in sections:
        story.append(Paragraph(sec_title, styles['CustomH']))
        items = [ListItem(Paragraph(b, styles['CustomBody'])) for b in bullets]
        story.append(ListFlowable(items, bulletType='bullet', start='circle', leftIndent=18))
        story.append(Spacer(1, 10))
    doc.build(story)


def main():
    build_internal_pptx(BASE / 'xbank_internal_finhouse_deck.pptx')
    build_client_pptx(BASE / 'xbank_client_deck.pptx')
    build_pdf(BASE / 'xbank_proje_analizi.pdf', 'X Bank Merchant Migration — Proje Analizi', [
        ('Yönetici Özeti', [
            'X Bank kendi sanal POS platformuna geçiyor; ilk 60 merchant taşındı, hedef yaklaşık 18.000 merchant.',
            'Bu dönüşüm teknik entegrasyondan çok yüksek hacimli bir merchant migration operasyonudur.',
            'Başarı için segmentasyon, operasyon disiplini, teknik enablement ve güçlü PMO birlikte gerekir.'
        ]),
        ('Kritik Gerçeklikler', [
            'Merchant’ların önemli kısmında teknik ekip yoktur.',
            'Hazır paket, custom ve ajans bağımlı entegrasyonlar birlikte bulunur.',
            'Toplu bildirim değil, aktif takip ve birebir yönlendirme gerekir.'
        ]),
        ('Önerilen Model', [
            'Migration factory yaklaşımı',
            'Wave bazlı rollout',
            'Pilot + playbook + hypercare modeli',
            'Hazır paketler için otomasyon ve geçiş kitleri'
        ])
    ])
    build_pdf(BASE / 'xbank_teklif_taslagi.pdf', 'X Bank Merchant Migration — Teklif Taslağı', [
        ('Kapsam', [
            'Merchant segmentasyonu ve geçiş metodolojisi',
            'Pilot faz planlama ve yürütme',
            'Merchant iletişim operasyonu',
            'Teknik onboarding koordinasyonu',
            'Raporlama ve hypercare'
        ]),
        ('Önerilen Ticari Çerçeve', [
            'Discovery / pilot için sabit başlangıç bedeli',
            'Merchant başı dönüşüm bedeli',
            'İşlem bazlı gelir modeli',
            'Gerekirse segment bazlı farklı ücretleme'
        ]),
        ('Neden Finhouse?', [
            'Payment domain bilgisi',
            'Saha gerçekliğini bilen migration yaklaşımı',
            'Operasyon + teknik enablement birleşik delivery modeli'
        ])
    ])

if __name__ == '__main__':
    main()
