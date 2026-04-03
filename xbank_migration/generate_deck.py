#!/usr/bin/env python3
"""
X Bank Merchant Migration — Detailed Client Deck Generator
Creates a professional 22-slide presentation using python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x1B, 0x2A, 0x4A)   # header / title bg
ACCENT     = RGBColor(0x1B, 0x8A, 0xC8)   # accent blue
LIGHT_BLUE = RGBColor(0xD6, 0xEA, 0xF8)   # box backgrounds
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)   # callout / highlight
GREEN      = RGBColor(0x1E, 0x8B, 0x4C)   # positive
RED        = RGBColor(0xC0, 0x39, 0x2B)   # risk/negative
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_GREY = RGBColor(0xF2, 0xF5, 0xF9)
MID_GREY   = RGBColor(0xBD, 0xC3, 0xC7)
FINHOUSE_ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(11), bold=False, color=DARK_TEXT,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_header(slide, title, subtitle=None):
    """Dark navy header bar across top"""
    add_rect(slide, 0, 0, 13.33, 1.05, fill_color=DARK_NAVY)
    add_textbox(slide, title, 0.3, 0.08, 10, 0.55,
                font_size=Pt(22), bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.62, 10, 0.38,
                    font_size=Pt(12), color=RGBColor(0xAA, 0xCC, 0xEE))
    # accent line
    add_rect(slide, 0, 1.05, 13.33, 0.04, fill_color=ACCENT)
    # footer bar
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_color=LIGHT_GREY)
    add_textbox(slide, "Finhouse — X Bank Merchant Migration Program  |  Gizli", 0.3, 7.18, 10, 0.28,
                font_size=Pt(8), color=MID_GREY)
    add_textbox(slide, "finhouse.ai", 11.8, 7.18, 1.3, 0.28,
                font_size=Pt(8), color=ACCENT, align=PP_ALIGN.RIGHT)


def bullet_block(slide, items, left, top, width, height,
                 title=None, title_color=DARK_NAVY, bullet="•",
                 font_size=Pt(11), bg_color=None, border_color=None):
    """A box with optional title and bullet items"""
    if bg_color:
        add_rect(slide, left, top, width, height, fill_color=bg_color, border_color=border_color, border_width=Pt(0.5))
    y = top + 0.1
    if title:
        add_textbox(slide, title, left+0.15, y, width-0.3, 0.32,
                    font_size=Pt(12), bold=True, color=title_color)
        y += 0.32
    for item in items:
        add_textbox(slide, f"{bullet}  {item}", left+0.15, y, width-0.3, 0.32,
                    font_size=font_size, color=DARK_TEXT)
        y += 0.30
    return y


def kpi_box(slide, value, label, left, top, w=2.2, h=1.1,
            val_color=ACCENT, bg=LIGHT_BLUE):
    add_rect(slide, left, top, w, h, fill_color=bg, border_color=ACCENT, border_width=Pt(0.75))
    add_textbox(slide, value, left+0.1, top+0.08, w-0.2, 0.5,
                font_size=Pt(22), bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, left+0.1, top+0.58, w-0.2, 0.45,
                font_size=Pt(9.5), color=DARK_TEXT, align=PP_ALIGN.CENTER)


def phase_box(slide, num, title, items, left, top, w=3.8, h=2.5, color=ACCENT):
    add_rect(slide, left, top, w, h, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, left, top, w, 0.38, fill_color=color)
    add_textbox(slide, f"FAZ {num}  {title}", left+0.1, top+0.05, w-0.2, 0.3,
                font_size=Pt(11), bold=True, color=WHITE)
    y = top + 0.46
    for item in items:
        add_textbox(slide, f"▸  {item}", left+0.12, y, w-0.25, 0.27,
                    font_size=Pt(10), color=DARK_TEXT)
        y += 0.27


def table_header_row(slide, headers, lefts, top, height=0.32, bg=DARK_NAVY):
    for i, (h, l) in enumerate(zip(headers, lefts)):
        w = lefts[i+1] - l if i+1 < len(lefts) else 13.0 - l
        add_rect(slide, l, top, w-0.02, height, fill_color=bg)
        add_textbox(slide, h, l+0.06, top+0.04, w-0.15, height-0.08,
                    font_size=Pt(9.5), bold=True, color=WHITE)


def table_data_row(slide, cells, lefts, top, height=0.28, bg=WHITE, border=MID_GREY):
    for i, (c, l) in enumerate(zip(cells, lefts)):
        w = lefts[i+1] - l if i+1 < len(lefts) else 13.0 - l
        add_rect(slide, l, top, w-0.02, height, fill_color=bg, border_color=border, border_width=Pt(0.3))
        add_textbox(slide, c, l+0.06, top+0.03, w-0.15, height-0.05,
                    font_size=Pt(9), color=DARK_TEXT)


# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
# Full dark navy background
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_NAVY)
# Accent strip
add_rect(slide, 0, 5.0, 13.33, 0.06, fill_color=ACCENT)
# Right-side accent panel
add_rect(slide, 9.8, 0, 3.53, 7.5, fill_color=RGBColor(0x12, 0x1E, 0x35))

add_textbox(slide, "X BANK", 0.6, 1.2, 9.0, 0.7,
            font_size=Pt(14), bold=True, color=ACCENT)
add_textbox(slide, "Sanal POS Merchant Migration Programı",
            0.6, 1.85, 9.0, 1.0,
            font_size=Pt(32), bold=True, color=WHITE)
add_textbox(slide, "Kontrollü • Ölçeklenebilir • Düşük Sürtünmeli Geçiş Modeli",
            0.6, 2.95, 9.0, 0.55,
            font_size=Pt(14), italic=True, color=RGBColor(0xAA, 0xCC, 0xEE))
# Divider
add_rect(slide, 0.6, 3.65, 5.0, 0.04, fill_color=ACCENT)

add_textbox(slide, "Hazırlayan:  Finhouse Danışmanlık", 0.6, 3.85, 5.0, 0.35,
            font_size=Pt(11), color=RGBColor(0xCC, 0xDD, 0xEE))
add_textbox(slide, "Tarih:  Mart 2026  |  Versiyon 2.0", 0.6, 4.2, 5.0, 0.35,
            font_size=Pt(11), color=RGBColor(0xCC, 0xDD, 0xEE))
add_textbox(slide, "GİZLİ — SADECE X BANK YETKİLİ KULLANICILARI İÇİN",
            0.6, 4.65, 9.0, 0.32, font_size=Pt(9),
            color=RGBColor(0x88, 0xAA, 0xCC), italic=True)

# Right panel content
add_textbox(slide, "18.000", 10.0, 1.6, 3.1, 0.75,
            font_size=Pt(40), bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(slide, "Merchant\nHedef Portföy", 10.0, 2.35, 3.1, 0.5,
            font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 10.2, 3.1, 2.7, 0.03, fill_color=ACCENT)
add_textbox(slide, "12 Ay", 10.0, 3.25, 3.1, 0.5,
            font_size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Program Süresi", 10.0, 3.75, 3.1, 0.35,
            font_size=Pt(11), color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
add_rect(slide, 10.2, 4.3, 2.7, 0.03, fill_color=ACCENT)
add_textbox(slide, "₺2+ Trilyon", 10.0, 4.45, 3.1, 0.5,
            font_size=Pt(20), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Türkiye e-Ticaret Pazarı\n(BKM 2024)", 10.0, 4.95, 3.1, 0.45,
            font_size=Pt(9.5), color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — YÖNETİCİ ÖZETİ
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Yönetici Özeti", "X Bank merchant migration programının temel çıktıları")
bg = add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# 4 key takeaway boxes
boxes = [
    ("🎯  Fırsat", DARK_NAVY,
     ["Payten'e ödenen dış maliyet içselleştirilir",
      "Merchant ilişkisi doğrudan sahiplenilir",
      "Cross-sell: BNPL, kredi, tahsilat"]),
    ("⚙️  Zorluk", RED,
     ["18.000 merchant heterojen teknik yapıda",
      "Toplu kampanya ile migration gerçekleşmiyor",
      "Aktif operasyon + teknik enablement şart"]),
    ("🏭  Yaklaşım", ACCENT,
     ["Migration Factory modeli",
      "3 fazlı: Discovery → Pilot → Scale",
      "Segment bazlı dalga yönetimi"]),
    ("📈  Beklenti", GREEN,
     ["12 ay içinde ≥%85 portföy dönüşümü",
      "Merchant churn <%15",
      "Brüt katkı marjı ~%30 (korumalı model)"]),
]
for i, (title, color, bullets) in enumerate(boxes):
    x = 0.25 + i * 3.25
    add_rect(slide, x, 1.25, 3.05, 2.55, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, x, 1.25, 3.05, 0.4, fill_color=color)
    add_textbox(slide, title, x+0.1, 1.28, 2.85, 0.34, font_size=Pt(11), bold=True, color=WHITE)
    y = 1.72
    for b in bullets:
        add_textbox(slide, f"▸  {b}", x+0.12, y, 2.82, 0.28, font_size=Pt(9.5), color=DARK_TEXT)
        y += 0.28

# KPI strip
add_rect(slide, 0.25, 4.05, 12.83, 0.04, fill_color=ACCENT)
add_textbox(slide, "Temel Metrikler", 0.25, 4.15, 4, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
kpis = [
    ("18.000", "Hedef Merchant"),
    ("12 Ay", "Program Süresi"),
    ("~27 Kişi", "Başlangıç Ekibi"),
    ("65 USD", "Birim Dönüşüm"),
    ("‰2", "İşlem Geliri"),
    ("~%85", "Hedef Dönüşüm"),
]
for i, (v, l) in enumerate(kpis):
    kpi_box(slide, v, l, 0.25 + i*2.15, 4.5, w=2.05, h=0.95)

# Finhouse proposition
add_rect(slide, 0.25, 5.6, 12.83, 0.95, fill_color=DARK_NAVY)
add_textbox(slide, "Finhouse Değer Önermesi:  QNB deneyimli ekip  ▪  Migration Factory metodolojisi  ▪  Operasyon + teknik enablement tek çatı  ▪  Ölçülebilir KPI taahhüdü",
            0.4, 5.75, 12.5, 0.65, font_size=Pt(11), color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — PAZAR BAĞLAMI
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Pazar Bağlamı", "Türkiye e-ticaret ve ödeme ekosistemi — stratejik zemin")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

add_textbox(slide, "Türkiye Ödeme Pazarı — Temel Göstergeler (BKM / TCMB 2024)",
            0.25, 1.18, 12.5, 0.38, font_size=Pt(13), bold=True, color=DARK_NAVY)

kpi_data = [
    ("₺2+ Trilyon", "Yıllık e-ticaret\nişlem hacmi"),
    ("100M+", "Bankalı kart\nkullanıcısı"),
    ("750M+", "Yıllık online\nkart işlem adedi"),
    ("%24", "YoY e-ticaret\nhacim büyümesi"),
    ("%99.9", "Hedef platform\nuptime"),
    ("T+1", "Hedef settlement\nsüresi"),
]
for i, (v, l) in enumerate(kpi_data):
    kpi_box(slide, v, l, 0.25 + i*2.15, 1.68, w=2.05, h=1.1)

# Ecosystem players
add_textbox(slide, "Sanal POS Ekosistemi — Temel Oyuncular",
            0.25, 2.95, 12.5, 0.35, font_size=Pt(12), bold=True, color=DARK_NAVY)

players = [
    ("Payten\n(Asseco SEE)", "Büyük bankaların\nSanal POS altyapısı\nPazar lideri", DARK_NAVY),
    ("iyzico\n(Mastercard)", "KOBİ + e-ticaret\nOdaklı tam çözüm\nDigital-first", ACCENT),
    ("Paynet\n(iyzico grubu)", "B2B & tedarikçi\nödeme altyapısı\nBayi ağları", GREEN),
    ("BKM Express\n/ Troy", "Ulusal kart\naltyapısı\nKamu destekli", ORANGE),
    ("X Bank\nYeni Platform", "Kendi geliştirdiği\nSanal POS\nMigration hedefinde", RED),
]
for i, (name, desc, color) in enumerate(players):
    x = 0.25 + i*2.55
    add_rect(slide, x, 3.35, 2.4, 1.5, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, x, 3.35, 2.4, 0.38, fill_color=color)
    add_textbox(slide, name, x+0.08, 3.37, 2.24, 0.34, font_size=Pt(9.5), bold=True, color=WHITE)
    add_textbox(slide, desc, x+0.1, 3.76, 2.2, 1.0, font_size=Pt(9), color=DARK_TEXT)

# Why now box
add_rect(slide, 0.25, 5.0, 12.83, 1.2, fill_color=RGBColor(0xE8, 0xF4, 0xFF), border_color=ACCENT, border_width=Pt(1))
add_textbox(slide, "💡  Neden Şimdi?", 0.4, 5.08, 3, 0.32, font_size=Pt(11), bold=True, color=DARK_NAVY)
reasons = ["Regülatör baskısı veri lokalizasyonunu zorunlu kılıyor",
           "Payten sözleşme yenileme dönemleri kritik pencere yaratıyor",
           "BNPL / açık bankacılık fırsatları merchant verisi sahipliği gerektiriyor"]
for i, r in enumerate(reasons):
    add_textbox(slide, f"▸  {r}", 0.4 + i*4.25, 5.48, 4.1, 0.62, font_size=Pt(10), color=DARK_TEXT)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — MEVCUT DURUM ANALİZİ
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Mevcut Durum Analizi", "X Bank — Payten altyapısından çıkışın stratejik gerekçesi")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

add_textbox(slide, "Payten Yapısının Maliyeti ve Kısıtları",
            0.25, 1.18, 6, 0.35, font_size=Pt(13), bold=True, color=DARK_NAVY)
add_textbox(slide, "Yeni Platformun Getirisi",
            6.8, 1.18, 6, 0.35, font_size=Pt(13), bold=True, color=GREEN)

# Left — current state
add_rect(slide, 0.25, 1.6, 6.3, 2.7, fill_color=WHITE, border_color=RED, border_width=Pt(1))
payten_items = [
    "Merchant başı: 80 USD / yıl işletim bedeli → dış maliyet",
    "İşlem başı: ‰3 komisyon → Payten'e ödeniyor",
    "Merchant ilişkisi Payten'e ait — cross-sell kör nokta",
    "Veri sahipliği yok — davranış içgörüsü erişilemiyor",
    "API kontrolü sınırlı — ürün geliştirme yavaşlıyor",
    "Fesih bildirimi: 6-12 ay önceden şart",
]
y = 1.68
for item in payten_items:
    add_textbox(slide, f"✗  {item}", 0.38, y, 6.05, 0.30, font_size=Pt(10), color=DARK_TEXT)
    y += 0.31

# Right — new state
add_rect(slide, 6.75, 1.6, 6.3, 2.7, fill_color=WHITE, border_color=GREEN, border_width=Pt(1))
new_items = [
    "Merchant başı: 65 USD dönüşüm geliri → Finhouse'a ödeme",
    "İşlem başı: ‰2 hizmet geliri → içselleşiyor",
    "Merchant ilişkisi doğrudan bankaya geçiyor",
    "Tam veri sahipliği: cross-sell, BNPL, kredi imkânı",
    "API tam kontrol: ürün geliştirme hız kazanıyor",
    "Migration tamamlandığında Payten sözleşmesi sonlanır",
]
y = 1.68
for item in new_items:
    add_textbox(slide, f"✓  {item}", 6.88, y, 6.05, 0.30, font_size=Pt(10), color=DARK_TEXT)
    y += 0.31

# Annual saving calculation
add_rect(slide, 0.25, 4.45, 12.83, 1.65, fill_color=DARK_NAVY)
add_textbox(slide, "Yıllık Tasarruf / Gelir Potansiyeli — 18.000 Merchant Tam Dönüşüm Senaryosu",
            0.4, 4.52, 12.4, 0.35, font_size=Pt(11), bold=True, color=WHITE)
calcs = [
    ("Payten işletim\nmaliyeti tasarrufu", "18.000 × $80", "= $1.44M / yıl", ACCENT),
    ("İşlem komisyonu\ntasarrufu (‰3→‰2)", "Hacim bağımlı", "≈ $500K–2M / yıl", GREEN),
    ("Setup / dönüşüm\ngeliri (tek seferlik)", "18.000 × $65", "= $1.17M", ORANGE),
    ("Cross-sell potansiyel\ngelir (tahmini)", "Merchant başı\n$150–500 / yıl", "= $2.7M–9M / yıl", RGBColor(0xCC, 0xAA, 0xFF)),
]
for i, (label, formula, result, color) in enumerate(calcs):
    x = 0.4 + i*3.1
    add_textbox(slide, label, x, 4.92, 2.7, 0.4, font_size=Pt(9), color=RGBColor(0xCC, 0xDD, 0xEE))
    add_textbox(slide, formula, x, 5.32, 2.7, 0.3, font_size=Pt(8.5), color=MID_GREY, italic=True)
    add_textbox(slide, result, x, 5.62, 2.7, 0.35, font_size=Pt(11), bold=True, color=color)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — STRATEJİK KAZANIM
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Stratejik Kazanım", "Kendi sanal POS platformuna geçişin çok boyutlu değeri")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

pillars = [
    ("💰  Maliyet\nİçselleştirme", DARK_NAVY,
     "Payten'e ödenen $80/merchant işletim bedeli + ‰3 komisyon bitiyor.\nToplam tasarruf potansiyeli: $2M+ / yıl"),
    ("🤝  Merchant\nİlişkisi", ACCENT,
     "Merchant verisi, tercihleri ve ödeme davranışı direkt X Bank'a ait oluyor.\nYukarı satış ve ürün ekleme için altın zemin."),
    ("🔗  Cross-Sell\nFırsatları", GREEN,
     "Merchant'a özel: İşletme kredisi, BNPL, maaş ödeme, tahsilat çözümü.\nOrtalama gelir potansiyeli: $150–500/merchant/yıl"),
    ("⚡  Platform\nKontrolü", ORANGE,
     "API yol haritası, güvenlik politikaları, yeni ürün lansmanları\nartık bankaya bağlı. İnovasyon hızı artıyor."),
]
for i, (title, color, desc) in enumerate(pillars):
    x = 0.25 + i*3.25
    add_rect(slide, x, 1.25, 3.05, 3.2, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, x, 1.25, 3.05, 0.55, fill_color=color)
    add_textbox(slide, title, x+0.1, 1.27, 2.85, 0.5, font_size=Pt(11), bold=True, color=WHITE)
    add_textbox(slide, desc, x+0.12, 1.88, 2.82, 2.45, font_size=Pt(10), color=DARK_TEXT)

# Bottom message
add_rect(slide, 0.25, 4.6, 12.83, 1.65, fill_color=DARK_NAVY)
add_textbox(slide, "Büyük Resim:  Merchant migration yalnızca teknik geçiş değil — X Bank için ödeme ekosisteminin kalıcı olarak içselleştirilmesidir.",
            0.45, 4.72, 12.3, 0.5, font_size=Pt(12), bold=True, color=WHITE)
add_textbox(slide, "Her geçen merchant, bankaya: (1) yıllık maliyet tasarrufu, (2) doğrudan ilişki, (3) veri varlığı, (4) cross-sell fırsatı getiriyor.\n"
            "Migration programının başarısı doğrudan X Bank'ın ödeme stratejisinin başarısıyla eşdeğerdir.",
            0.45, 5.28, 12.3, 0.85, font_size=Pt(10), color=RGBColor(0xCC, 0xDD, 0xEE))

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — GLOBAL BENCHMARK
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Global Benchmark", "Stripe Connect, Adyen ve sektör liderleri — kanıtlanmış migration modelleri")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# Stripe box
add_rect(slide, 0.25, 1.2, 5.9, 2.8, fill_color=WHITE, border_color=ACCENT, border_width=Pt(1.5))
add_rect(slide, 0.25, 1.2, 5.9, 0.4, fill_color=ACCENT)
add_textbox(slide, "Stripe Connect — Platform Migration Modeli", 0.38, 1.23, 5.6, 0.34, font_size=Pt(11), bold=True, color=WHITE)
stripe_items = [
    "Standard → Express → Custom olgunluk skalası",
    "Incremental migration: batch'lerle kademeli geçiş",
    "Prefill yaklaşımı: önceki veriyle sürtünme azaltımı",
    "Platform (banka) KYC/MASAK sorumluluğunu üstlenir",
    "Backward compatibility: eski endpoint 18 ay canlı",
    "Fiyat: %2.9 + $0.30 (standart), IC+ (kurumsal)",
]
y = 1.68
for item in stripe_items:
    add_textbox(slide, f"▸  {item}", 0.38, y, 5.65, 0.27, font_size=Pt(9.5), color=DARK_TEXT)
    y += 0.28

# Adyen box
add_rect(slide, 6.4, 1.2, 5.9, 2.8, fill_color=WHITE, border_color=GREEN, border_width=Pt(1.5))
add_rect(slide, 6.4, 1.2, 5.9, 0.4, fill_color=GREEN)
add_textbox(slide, "Adyen for Platforms — Büyük Ölçekli Onboarding", 6.53, 1.23, 5.6, 0.34, font_size=Pt(11), bold=True, color=WHITE)
adyen_items = [
    "İşlem başı sabit: $0.13 (~₺5) processing fee",
    "Interchange++ modeli: kart tipine göre %1.3–4.5",
    "Account Holder modeli: KYC platform sorumluluğu",
    "Farklı ülkeler için ayrı ödeme yöntemi anlaşması",
    "Risk yönetimi: platform (banka) chargeback üstlenir",
    "Türkiye: BDDK yetkili banka direkt KYC yürütür",
]
y = 1.68
for item in adyen_items:
    add_textbox(slide, f"▸  {item}", 6.53, y, 5.65, 0.27, font_size=Pt(9.5), color=DARK_TEXT)
    y += 0.28

# Benchmark table
add_textbox(slide, "Büyük Ölçekli Migration Referans Projeleri (Sektör Verisi)",
            0.25, 4.1, 12.5, 0.32, font_size=Pt(11), bold=True, color=DARK_NAVY)
hdrs = ["Proje", "Ölçek", "Süre", "Anahtar Ders"]
lefts = [0.25, 4.5, 7.0, 8.8]
table_header_row(slide, hdrs, lefts, 4.48)
rows = [
    ("Worldpay / FIS konsolidasyon", "500K+ merchant", "18–24 ay", "Segment bazlı kademeli geçiş zorunlu"),
    ("PayPal / Braintree platform", "100K+ platform", "12 ay", "API uyumu korundu, eski endpoint 18 ay canlı"),
    ("Fiserv / First Data", "~6M merchant", "36 ay", "\"Never dark\" yaklaşımı — parallel run şart"),
    ("Türkiye yerel banka (Payten→kendi)", "~15K merchant", "14 ay", "KOBİ en yavaş; e-ticaret paket en hızlı segment"),
]
for i, row in enumerate(rows):
    bg = WHITE if i % 2 == 0 else LIGHT_GREY
    table_data_row(slide, row, lefts, 4.82 + i*0.3, bg=bg)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — MERCHANT SEGMENTASYONu
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Merchant Segmentasyonu", "18.000 merchant homojen değil — segment bazlı yaklaşım zorunlu")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

segments = [
    ("A", "Hazır E-Ticaret\nPaketleri", "%40 — ~7.200 merchant",
     ["Ticimax, IdeaSoft, T-Soft, İkas, WooCommerce",
      "Plugin/adaptör ile toplu enablement mümkün",
      "Migration sürtünmesi en düşük segment",
      "İş yükü çarpanı: 0.5x — en verimli"],
     GREEN, "~30 USD/merchant"),
    ("B", "Standart API\nEntegrasyonu", "%35 — ~6.300 merchant",
     ["Kendi IT ekibi olan orta/büyük merchant'lar",
      "API/callback/3D/hash/sertifika koordinasyonu",
      "Sandbox + test ortamı + go-live checklist",
      "İş yükü çarpanı: 1.0x — baz segment"],
     ACCENT, "~61 USD/merchant"),
    ("C", "Ajans / Dış Yazılım\nBağımlı", "%20 — ~3.600 merchant",
     ["Merchant yerine dış firmaya iletişim",
      "Koordinasyon zinciri uzuyor",
      "Ajans motivasyonu ve uyum kritik faktör",
      "İş yükü çarpanı: 2.5x — en yüksek maliyet"],
     ORANGE, "~153 USD/merchant"),
    ("D", "Pasif / Düşük\nHacimli KOBİ", "%5 — ~900 merchant",
     ["Teknik kapasitesi en düşük grup",
      "Adım adım rehberlik, ekran paylaşımı",
      "Ticari öncelik düşük — wave sonuna bırakılır",
      "İş yükü çarpanı: 0.8x — hafif operasyon"],
     MID_GREY, "~49 USD/merchant"),
]
for i, (letter, name, count, bullets, color, cost) in enumerate(segments):
    x = 0.25 + i*3.25
    add_rect(slide, x, 1.25, 3.05, 3.4, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, x, 1.25, 3.05, 0.55, fill_color=color)
    add_textbox(slide, f"Seg. {letter}  —  {name}", x+0.1, 1.27, 2.85, 0.5, font_size=Pt(10), bold=True, color=WHITE)
    add_textbox(slide, count, x+0.1, 1.82, 2.85, 0.28, font_size=Pt(9), bold=True, color=color)
    y = 2.15
    for b in bullets:
        add_textbox(slide, f"▸  {b}", x+0.12, y, 2.82, 0.27, font_size=Pt(9), color=DARK_TEXT)
        y += 0.27
    add_rect(slide, x, 4.55, 3.05, 0.35, fill_color=color)
    add_textbox(slide, f"Birim maliyet: {cost}", x+0.1, 4.6, 2.85, 0.27, font_size=Pt(9.5), bold=True, color=WHITE)

# Migration priority arrow
add_rect(slide, 0.25, 4.98, 12.83, 0.04, fill_color=ACCENT)
add_textbox(slide, "Migration Öncelik Sırası  →  A Segmenti önce (en hızlı dönüşüm) → B → D → C (en fazla destek)",
            0.25, 5.05, 12.5, 0.32, font_size=Pt(10), color=DARK_NAVY, bold=True)
add_textbox(slide, "Önemli Not: Segment bazlı fiyatlama teklife yansıtılmalı — Custom (C) merchant'lar standart ücretin 2x'i olarak ayrıca fiyatlanır.",
            0.25, 5.42, 12.5, 0.32, font_size=Pt(10), color=RED, italic=True)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — 3 AŞAMALI YAKLAŞIM
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "3 Aşamalı Yaklaşım", "Discovery → Pilot → Scale — kontrollü ve ölçülebilir ilerleme modeli")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# Phase timeline bar
for i, (label, color, weeks) in enumerate([
    ("FAZ 1  Discovery & Pilot  (Hafta 1–8)", DARK_NAVY, 3.0),
    ("FAZ 2  Endüstriyel Migrasyon  (Ay 3–11)", ACCENT, 6.5),
    ("FAZ 3  Hypercare & Kapanış  (Ay 12)", GREEN, 2.5),
]):
    x = 0.25 + sum([3.0, 6.5, 2.5][:i])
    add_rect(slide, x, 1.18, weeks, 0.5, fill_color=color)
    add_textbox(slide, label, x+0.08, 1.22, weeks-0.15, 0.4, font_size=Pt(9.5), bold=True, color=WHITE)

# Arrow connector
add_textbox(slide, "→", 3.25, 1.25, 0.3, 0.35, font_size=Pt(16), bold=True, color=WHITE)
add_textbox(slide, "→", 9.75, 1.25, 0.3, 0.35, font_size=Pt(16), bold=True, color=WHITE)

# Phase detail boxes
phase_box(slide, "1", "DISCOVERY & PILOT", [
    "Merchant envanteri alınır, segmentasyon yapılır",
    "Pilot kapsamı: 100–300 merchant",
    "Arama scriptleri & teknik playbook oluşturulur",
    "Gerçek efor ve dönüşüm süresi ölçülür",
    "Çıktı: Segment stratejisi + ekip planı",
    "Süre: 8 hafta  |  Sabit bedel: $150K–250K",
], 0.25, 1.82, w=4.0, h=2.8, color=DARK_NAVY)

phase_box(slide, "2", "ENDÜSTRİYEL MİGRASYON", [
    "Dalga bazlı yayılım: A → B → D → C",
    "Merchant iletişim + teknik onboarding",
    "Test ortamı + canlıya geçiş koordinasyonu",
    "Haftalık/aylık KPI raporlaması",
    "Issue yönetimi + banka koordinasyonu",
    "Hedef: Aylık 1.500–2.000 merchant",
], 4.4, 1.82, w=4.5, h=2.8, color=ACCENT)

phase_box(slide, "3", "HYPERCARE & KAPANIŞ", [
    "Geçiş sonrası ilk 30 gün yoğun destek",
    "Açık issue takibi ve çözümü",
    "KPI final raporu hazırlanır",
    "Operasyon devir teslimi (banka ekibine)",
    "Retrospective & öğrenim belgesi",
    "Süre: 4–6 hafta",
], 9.05, 1.82, w=4.0, h=2.8, color=GREEN)

# Key success factors
add_textbox(slide, "Başarı İçin Kritik Unsurlar", 0.25, 4.75, 12.5, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
factors = [
    ("📊  Doğru Segmentasyon", "Pilot öncesi tamamlanmış\nmust-have"),
    ("🧪  Pilot Öğrenimleri", "Wave 1 çıktıları\nsonraki dalgaları şekillendirir"),
    ("🔌  Plugin/Adaptör", "Hazır paket merchant'lar\niçin self-serve enablement"),
    ("📋  Net Başarı Kriteri", "Canlı ilk başarılı işlem\n= ödeme tetikleyicisi"),
    ("🛡  Parallel Run", "Büyük merchant'lar için\neski+yeni sistem eşzamanlı"),
]
for i, (title, desc) in enumerate(factors):
    x = 0.25 + i*2.55
    add_rect(slide, x, 5.12, 2.4, 1.1, fill_color=LIGHT_BLUE, border_color=ACCENT, border_width=Pt(0.5))
    add_textbox(slide, title, x+0.1, 5.18, 2.2, 0.32, font_size=Pt(9.5), bold=True, color=DARK_NAVY)
    add_textbox(slide, desc, x+0.1, 5.5, 2.2, 0.65, font_size=Pt(9), color=DARK_TEXT)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — DISCOVERY FAZI DETAYI
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Faz 1 — Discovery & Pilot Detayı", "Hafta 1–8  |  Sabit bedel: $150.000–250.000")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# Week breakdown
add_textbox(slide, "Haftalık Çalışma Planı", 0.25, 1.18, 7, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
week_data = [
    ("H1–2", "Envanter & Analiz", "Merchant master data alınır, segment dağılımı çıkarılır, önceliklendirme yapılır"),
    ("H3–4", "Süreç Tasarımı", "Migration playbook, iletişim scriptleri, teknik onboarding checklist hazırlanır"),
    ("H5–6", "Pilot Yürütme", "100–300 merchant pilot başlatılır, segment A ve B öncelikli"),
    ("H7–8", "Öğrenim & Plan", "Pilot çıktıları analiz edilir, wave planı güncellenir, kapasite netleştirilir"),
]
y = 1.58
for week, title, desc in week_data:
    add_rect(slide, 0.25, y, 0.9, 0.58, fill_color=DARK_NAVY)
    add_textbox(slide, week, 0.25, y+0.12, 0.9, 0.35, font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.2, y, 6.3, 0.58, fill_color=WHITE, border_color=ACCENT, border_width=Pt(0.5))
    add_textbox(slide, title, 1.32, y+0.04, 2.2, 0.28, font_size=Pt(10), bold=True, color=DARK_NAVY)
    add_textbox(slide, desc, 3.6, y+0.04, 3.8, 0.5, font_size=Pt(9), color=DARK_TEXT)
    y += 0.63

# X Bank'tan istenen veri
add_rect(slide, 7.65, 1.18, 5.4, 3.45, fill_color=WHITE, border_color=ORANGE, border_width=Pt(1.5))
add_rect(slide, 7.65, 1.18, 5.4, 0.38, fill_color=ORANGE)
add_textbox(slide, "X Bank'tan İstenen Başlangıç Verisi", 7.78, 1.21, 5.1, 0.32, font_size=Pt(11), bold=True, color=WHITE)
data_items = [
    "Merchant master data (ID, ad, iletişim, segment)",
    "Son 12 ay işlem hacmi ve adedi (özet)",
    "Platform / entegrasyon tipi dağılımı",
    "İlk 60 merchant geçişinden çıkan öğrenimler",
    "Teknik dokümantasyon & sandbox durumu",
    "Bankanın iç koordinatör atanması",
    "Hedef takvim ve başarı kriteri tanımı",
]
y = 1.62
for item in data_items:
    add_textbox(slide, f"□  {item}", 7.78, y, 5.1, 0.27, font_size=Pt(9.5), color=DARK_TEXT)
    y += 0.29

# Discovery outputs
add_textbox(slide, "Discovery Çıktıları", 0.25, 4.2, 12.5, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
outputs = [
    ("Segment Stratejisi", "Her segment için ayrı\nmigration playbook"),
    ("Pilot Raporu", "Dönüşüm süresi, engeller,\nöğrenimler"),
    ("Ekip & Kapasite Planı", "Faz 2 için ölçeklenmiş\nekip yapısı"),
    ("Risk Matrisi", "Tespit edilen riskler\nve azaltma planı"),
    ("Nihai Teklif", "Faz 2–3 için güncellenmiş\nSOW ve fiyat"),
]
for i, (title, desc) in enumerate(outputs):
    x = 0.25 + i*2.55
    add_rect(slide, x, 4.58, 2.4, 1.15, fill_color=LIGHT_BLUE, border_color=DARK_NAVY, border_width=Pt(0.5))
    add_textbox(slide, title, x+0.1, 4.64, 2.2, 0.3, font_size=Pt(9.5), bold=True, color=DARK_NAVY)
    add_textbox(slide, desc, x+0.1, 4.95, 2.2, 0.7, font_size=Pt(9), color=DARK_TEXT)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — PILOT FAZI DETAYI
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Faz 2 — Endüstriyel Migrasyon: Dalga Bazlı Yayılım", "Ay 3–11  |  Aylık hedef: 1.500–2.000 merchant")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# Wave plan
add_textbox(slide, "Wave Planı — 9 Aylık Migrasyon Döngüsü", 0.25, 1.18, 9, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
waves = [
    ("Wave 1\n(Ay 3–4)", "Segment A\nHazır Paket", "~3.000", "Plugin adaptörleri\naktive edilir", ACCENT),
    ("Wave 2\n(Ay 4–6)", "Segment B\nStandart API", "~3.200", "Teknik onboarding\nyoğunlaşır", GREEN),
    ("Wave 3\n(Ay 6–8)", "Segment A+B\nKalan", "~4.000", "Self-serve kanal\ndevrede", DARK_NAVY),
    ("Wave 4\n(Ay 7–9)", "Segment D\nKOBİ", "~900", "Hafif operasyon\nmodeli", ORANGE),
    ("Wave 5\n(Ay 9–11)", "Segment C\nCustom/Ajans", "~6.900", "En yoğun teknik\ndestek dönemi", RED),
]
for i, (wave, segment, count, note, color) in enumerate(waves):
    x = 0.25 + i*2.55
    add_rect(slide, x, 1.57, 2.4, 2.2, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, x, 1.57, 2.4, 0.42, fill_color=color)
    add_textbox(slide, wave, x+0.1, 1.59, 2.2, 0.38, font_size=Pt(9.5), bold=True, color=WHITE)
    add_textbox(slide, segment, x+0.1, 2.05, 2.2, 0.32, font_size=Pt(9.5), bold=True, color=color)
    add_textbox(slide, f"~{count} merchant", x+0.1, 2.38, 2.2, 0.27, font_size=Pt(11), bold=True, color=DARK_NAVY)
    add_textbox(slide, note, x+0.1, 2.68, 2.2, 0.45, font_size=Pt(9), color=DARK_TEXT)

# Migration process steps
add_textbox(slide, "Her Merchant İçin Süreç Akışı", 0.25, 3.9, 12.5, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
steps = ["1. Envanter\n& Temas", "2. Teknik\nReadiness", "3. Test\nOrtamı", "4. Canlı\nGeçiş", "5. Hypercare\n(30 gün)"]
step_colors = [DARK_NAVY, ACCENT, ORANGE, GREEN, RGBColor(0x7D, 0x3C, 0x98)]
for i, (step, color) in enumerate(zip(steps, step_colors)):
    x = 0.25 + i*2.55
    add_rect(slide, x, 4.28, 2.4, 0.65, fill_color=color)
    add_textbox(slide, step, x+0.1, 4.32, 2.2, 0.55, font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_textbox(slide, "→", x+2.4, 4.4, 0.18, 0.35, font_size=Pt(14), bold=True, color=DARK_NAVY)

# KPI monitoring
add_rect(slide, 0.25, 5.1, 12.83, 1.0, fill_color=DARK_NAVY)
add_textbox(slide, "Haftalık KPI İzleme", 0.4, 5.15, 4, 0.28, font_size=Pt(10), bold=True, color=WHITE)
kpi_items = ["Temas kurulan / Toplam", "Test tamamlanan", "Canlıya geçen", "Açık issue", "Ortalama geçiş süresi"]
for i, kpi in enumerate(kpi_items):
    add_textbox(slide, f"▸ {kpi}", 0.4 + (i % 3)*4.2, 5.45 + (i//3)*0.28, 4.0, 0.28, font_size=Pt(9.5), color=RGBColor(0xCC, 0xDD, 0xEE))

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — OPERASYON MODELİ
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Operasyon Modeli", "Migration Factory — Ekip Yapısı ve Roller")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# Org chart style boxes
add_rect(slide, 5.2, 1.22, 2.9, 0.58, fill_color=DARK_NAVY)
add_textbox(slide, "Program Direktörü\n(Program Manager)", 5.25, 1.25, 2.8, 0.52, font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(slide, 6.65, 1.8, 0.04, 0.4, fill_color=DARK_NAVY)

# Row 2
row2 = [
    (0.5, "Project Lead\nStream A", ACCENT),
    (3.65, "Project Lead\nStream B", ACCENT),
    (6.8, "Solution\nArchitect", ORANGE),
    (10.0, "Sr. Payment\nConsultant", ORANGE),
]
for x, label, color in row2:
    add_rect(slide, x, 2.22, 2.5, 0.52, fill_color=color)
    add_textbox(slide, label, x+0.08, 2.24, 2.35, 0.48, font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Row 3 — execution teams
add_textbox(slide, "Yürütme Ekipleri", 0.25, 2.95, 12.5, 0.3, font_size=Pt(11), bold=True, color=DARK_NAVY)

teams = [
    ("Merchant\nOperasyon", "10 Migration Specialist\nArama, takip, onboarding\nrandevu & checklist yönetimi", GREEN, 5.0),
    ("Teknik\nOnboarding", "5 Technical Spec.\nAPI entegrasyon, test\nhata analizi, troubleshoot", ACCENT, 5.0),
    ("PMO &\nRaporlama", "2 PMO Analyst\nKPI dashboard, wave plan\nhaftalık yönetim raporu", DARK_NAVY, 3.5),
    ("QA & Kalite", "2 QA Specialist\nTest kontrol, playbook\nkalite güvence", ORANGE, 3.0),
    ("Support &\nHypercare", "3 Support Spec.\nPost-live destek, L1/L2\nticket yönetimi", RGBColor(0x7D, 0x3C, 0x98), 3.5),
]
x = 0.25
for title, desc, color, width in teams:
    add_rect(slide, x, 3.3, width, 1.5, fill_color=WHITE, border_color=color, border_width=Pt(1.2))
    add_rect(slide, x, 3.3, width, 0.38, fill_color=color)
    add_textbox(slide, title, x+0.1, 3.33, width-0.2, 0.33, font_size=Pt(9.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x+0.1, 3.72, width-0.2, 1.0, font_size=Pt(9), color=DARK_TEXT)
    x += width + 0.25

# Headcount summary
add_rect(slide, 0.25, 5.0, 12.83, 1.15, fill_color=DARK_NAVY)
add_textbox(slide, "Ekip Özeti — Baz Senaryo (Faz 2 Yoğun Dönem)",
            0.4, 5.06, 12.3, 0.3, font_size=Pt(10), bold=True, color=WHITE)
hc_items = [
    ("1", "Program Mgr"), ("2", "Project Lead"), ("1", "Solution Arch"),
    ("1", "Sr. Consultant"), ("5", "Tech Spec."), ("10", "Migration Spec."),
    ("2", "PMO Analyst"), ("2", "QA"), ("3", "Support"), ("27", "Toplam"),
]
for i, (n, label) in enumerate(hc_items):
    x = 0.4 + i*1.28
    color = ACCENT if label != "Toplam" else ORANGE
    add_textbox(slide, n, x, 5.42, 1.1, 0.32, font_size=Pt(14), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label, x, 5.72, 1.1, 0.28, font_size=Pt(8), color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — TEKNOLOJİ & ARAÇLAR
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Teknoloji & Araçlar", "Migration Factory'nin teknik omurgası")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

tool_cats = [
    ("Proje & İş Yönetimi", DARK_NAVY, [
        "Jira / Asana — wave planı, issue takibi",
        "Confluence / Notion — playbook & dokümantasyon",
        "Microsoft Teams — banka koordinasyonu",
        "Power BI / Looker — KPI dashboard",
    ]),
    ("Merchant İletişim", ACCENT, [
        "CRM (HubSpot / Salesforce) — merchant takibi",
        "SMS gateway — toplu bildirim kampanyaları",
        "E-posta otomasyon — onboarding dizisi",
        "Çağrı merkezi entegrasyonu — outbound dialing",
    ]),
    ("Teknik Test & Entegrasyon", GREEN, [
        "Postman / Insomnia — API test ve dokümantasyon",
        "Sandbox ortamı — her segment için izole test",
        "Webhook simulator — callback doğrulama",
        "PCI-DSS SAQ-A / SAQ-A-EP compliance araçları",
    ]),
    ("E-Ticaret Platform Adaptörleri", ORANGE, [
        "WooCommerce plugin — otomasyon geçiş paketi",
        "Shopify App — hazır entegrasyon",
        "Ticimax / IdeaSoft / T-Soft modülleri",
        "OpenCart, PrestaShop adaptörleri",
    ]),
]
for i, (title, color, items) in enumerate(tool_cats):
    x = 0.25 + (i % 2)*6.55
    y = 1.22 + (i // 2)*2.7
    add_rect(slide, x, y, 6.3, 2.55, fill_color=WHITE, border_color=color, border_width=Pt(1.2))
    add_rect(slide, x, y, 6.3, 0.38, fill_color=color)
    add_textbox(slide, title, x+0.12, y+0.04, 6.05, 0.3, font_size=Pt(11), bold=True, color=WHITE)
    yy = y + 0.44
    for item in items:
        add_textbox(slide, f"▸  {item}", x+0.15, yy, 5.95, 0.27, font_size=Pt(9.5), color=DARK_TEXT)
        yy += 0.28

add_rect(slide, 0.25, 6.68, 12.83, 0.04, fill_color=ACCENT)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — RİSK MATRİSİ
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Risk Matrisi", "Önceden tespit edilmiş riskler ve azaltma stratejileri")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

hdrs = ["Risk", "Olasılık", "Etki", "Azaltma Stratejisi"]
lefts = [0.25, 5.3, 7.0, 8.8]
table_header_row(slide, hdrs, lefts, 1.22)

risk_data = [
    ("Merchant churn — platforma geçmeme", "Yüksek", "Yüksek", "Teşvik paketi, SLA taahhüdü, eski sistem paralel canlı"),
    ("Teknik entegrasyon gecikmesi", "Orta", "Yüksek", "Sandbox ortamı, plugin katalog, API backward compat"),
    ("Custom merchant oranı %15 aşıyor", "Orta", "Yüksek", "Segment bazlı fiyatlama, discovery fazında tespit"),
    ("KYC/MASAK ret oranı yüksek", "Orta", "Orta", "Başvuru öncesi ön eleme, e-devlet entegrasyonu"),
    ("Banka veri kalitesi düşük", "Orta", "Orta", "Discovery fazında doğrulama, temizleme protokolü"),
    ("Proje 12 ay aşıyor", "Düşük–Orta", "Orta", "Aylık uzama bedeli maddesi sözleşmede yer alır"),
    ("Kur hareketi (TL güçlenmesi)", "Orta", "Yüksek", "Sözleşmede kur güvence maddesi ±%15 bandı"),
    ("Platform uptime sorunu", "Düşük", "Yüksek", "%99.9 SLA taahhüdü + monitoring + DR planı"),
    ("Ekip kapasitesi yetersizliği", "Düşük", "Orta", "Kademeli ramp-up, faz bazlı ölçekleme"),
]
row_colors = [
    (RED, RGBColor(0xFF, 0xEB, 0xEE)),   # H/H — red bg light
    (ORANGE, RGBColor(0xFF, 0xF3, 0xE0)),
    (ORANGE, RGBColor(0xFF, 0xF3, 0xE0)),
    (ORANGE, RGBColor(0xFF, 0xF3, 0xE0)),
    (ORANGE, RGBColor(0xFF, 0xF3, 0xE0)),
    (MID_GREY, WHITE),
    (ORANGE, RGBColor(0xFF, 0xF3, 0xE0)),
    (MID_GREY, WHITE),
    (MID_GREY, WHITE),
]
for i, (row, (text_c, bg_c)) in enumerate(zip(risk_data, row_colors)):
    table_data_row(slide, row, lefts, 1.57 + i*0.52, height=0.5, bg=bg_c, border=MID_GREY)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — MALİYET MODELİ
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Maliyet Modeli", "Faz bazlı ekip yapısı ve operasyon maliyetleri — baz senaryo")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

add_textbox(slide, "Ekip Maliyet Tablosu — Faz Bazlı (Fully Loaded TL/ay)",
            0.25, 1.18, 12.5, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
hdrs = ["Rol", "Faz 1 (Ay 1–2)", "Faz 2 (Ay 3–11)", "Faz 3 (Ay 12)", "Birim FL TL/ay"]
lefts2 = [0.25, 4.5, 6.6, 8.7, 10.8]
table_header_row(slide, hdrs, lefts2, 1.55, height=0.3)
team_rows = [
    ("Program Manager", "1", "1", "1", "240.000"),
    ("Project / Stream Lead", "1", "2", "1", "175.000"),
    ("Solution Architect", "1", "1", "0,5", "210.000"),
    ("Sr. Payment Consultant", "1", "1", "0,5", "195.000"),
    ("Technical Integration Spec.", "3", "5", "2", "128.000"),
    ("Merchant Migration Spec.", "2", "9", "4", "88.000"),
    ("PMO / Reporting", "1", "2", "1", "95.000"),
    ("QA / Process Control", "1", "2", "1", "92.000"),
    ("Support / Hypercare", "1", "3", "2", "83.000"),
    ("Aylık Toplam (TL)", "1.650.000", "3.050.000", "1.578.500", "—"),
]
for i, row in enumerate(team_rows):
    bg = LIGHT_BLUE if row[0] == "Aylık Toplam (TL)" else (WHITE if i % 2 == 0 else LIGHT_GREY)
    bold_flag = row[0] == "Aylık Toplam (TL)"
    table_data_row(slide, row, lefts2, 1.88 + i*0.3, height=0.28, bg=bg)

# Summary boxes
add_textbox(slide, "Senaryo Özeti", 0.25, 4.98, 12.5, 0.3, font_size=Pt(11), bold=True, color=DARK_NAVY)
scenarios = [
    ("Agresif\n(Optimize)", "~₺34M\n~$49/merchant", "Yüksek otomasyon,\ntemiz veri, banka desteği", GREEN),
    ("Baz Senaryo\n(Önerilen)", "~₺41.7M\n~$61/merchant", "Makul koşullar,\nstandart ekip yapısı", ACCENT),
    ("Koruyucu\n(Risk Tamponlu)", "~₺52.5M\n~$77/merchant", "Veri sorunu, kur\nriski, proje uzaması", ORANGE),
]
for i, (label, cost, cond, color) in enumerate(scenarios):
    x = 0.25 + i*4.25
    add_rect(slide, x, 5.33, 4.0, 1.3, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, x, 5.33, 4.0, 0.4, fill_color=color)
    add_textbox(slide, label, x+0.1, 5.35, 3.8, 0.35, font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, cost, x+0.1, 5.78, 3.8, 0.4, font_size=Pt(13), bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, cond, x+0.1, 6.18, 3.8, 0.38, font_size=Pt(9), color=DARK_TEXT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — GELİR SENARYOLARI
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Gelir Senaryoları", "3 senaryo analizi — setup fee + işlem geliri birlikte değerlendirmesi")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

add_textbox(slide, "Önerilen Teklif Yapısı",
            0.25, 1.18, 12.5, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
fee_components = [
    ("Keşif & Pilot\nSabit Bedeli", "$150K–250K\n(sabit, 2 ay)", DARK_NAVY),
    ("Merchant Başı\nMigrasyon (Standart)", "65 USD\nSeg. A & B", ACCENT),
    ("Merchant Başı\nMigrasyon (Custom)", "120–150 USD\nSeg. C (ajans)", ORANGE),
    ("İşlem Bazlı\nGelir", "‰1.5–2.0\nAy 3'ten itibaren", GREEN),
]
for i, (label, value, color) in enumerate(fee_components):
    x = 0.25 + i*3.25
    add_rect(slide, x, 1.56, 3.05, 0.9, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_textbox(slide, label, x+0.1, 1.6, 2.85, 0.35, font_size=Pt(9.5), bold=True, color=color)
    add_textbox(slide, value, x+0.1, 1.95, 2.85, 0.45, font_size=Pt(11), bold=True, color=DARK_TEXT)

# 3-scenario P&L
add_textbox(slide, "Gelir-Gider Analizi — 3 Senaryo", 0.25, 2.63, 12.5, 0.3, font_size=Pt(12), bold=True, color=DARK_NAVY)
hdrs = ["Kalem", "Senaryo A\n(Korumalı)", "Senaryo B\n(Sadece Setup)", "Senaryo C\n(Optimize)"]
lefts3 = [0.25, 4.8, 7.4, 10.0]
table_header_row(slide, hdrs, lefts3, 2.98)
pl_rows = [
    ("Discovery + Pilot Sabit Gelir", "₺6.0M", "—", "₺6.0M"),
    ("Setup Fee Geliri (15K merchant × ₺2.470)", "₺37.1M", "₺44.5M", "₺37.1M"),
    ("İşlem Geliri (ihtiyatlı, yıllık)", "₺24.0M", "—", "₺32.0M"),
    ("TOPLAM GELİR", "₺67.1M", "₺44.5M", "₺75.1M"),
    ("Toplam Operasyon Maliyeti", "(₺47.0M)", "(₺47.0M)", "(₺41.0M)"),
    ("BRÜT KATKI", "₺20.1M", "(₺2.5M)", "₺34.1M"),
    ("Brüt Katkı Marjı", "%30", "Negatif", "%45"),
]
pl_highlights = [False, False, False, True, False, True, True]
for i, (row, hl) in enumerate(zip(pl_rows, pl_highlights)):
    bg = LIGHT_BLUE if hl else (WHITE if i % 2 == 0 else LIGHT_GREY)
    table_data_row(slide, row, lefts3, 3.32 + i*0.3, height=0.28, bg=bg)

# Key insight
add_rect(slide, 0.25, 5.52, 12.83, 0.95, fill_color=DARK_NAVY)
add_textbox(slide, "🔑  Stratejik Sonuç:  Yalnızca setup fee ile çalışmak (Senaryo B) riskli ve negatif marjlı olabilir.\n"
            "İşlem gelirinin sözleşmeye eklenmesi stratejik zorunluluktur. Korumalı model (%30 marj) önerilen pozisyon.",
            0.4, 5.6, 12.3, 0.8, font_size=Pt(10), bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — BAŞARI KPI'LARI
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Başarı KPI'ları & Ölçüm Çerçevesi", "Ölçülebilir taahhütler — migration programının yönetim aracı")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

hdrs = ["KPI Metriği", "Hedef", "Ölçüm Frekansı", "Sınır Değer"]
lefts4 = [0.25, 5.5, 8.4, 10.8]
table_header_row(slide, hdrs, lefts4, 1.22)
kpi_rows = [
    ("Migration tamamlanma oranı", "≥%85 portföy, 12 ay içinde", "Aylık", "<%70 → eskalasyon"),
    ("Merchant churn oranı (migration kaynaklı)", "<%15", "Aylık", ">%20 → kriz protokolü"),
    ("Onboarding süresi — KOBİ segment", "≤3 iş günü", "Haftalık", ">5 gün → süreç revizyonu"),
    ("Onboarding süresi — Kurumsal segment", "≤10 iş günü", "Haftalık", ">15 gün → dedicated kaynak"),
    ("Platform uptime (X Bank sanal POS)", "≥%99.9", "Günlük", "<%99.5 → SLA ihlali"),
    ("İşlem başarı oranı (3DS dahil)", "≥%97", "Günlük", "<%95 → teknik analiz"),
    ("Chargeback oranı", "<%0.5", "Aylık", ">%1 → risk protokolü (Visa/MC limit)"),
    ("Merchant NPS", "≥40", "Çeyreklik", "<30 → müşteri deneyimi revizyonu"),
    ("Teknik destek ticket çözüm süresi", "≤48 saat (L1/L2)", "Haftalık", ">72 saat → kaynak artışı"),
    ("İlk 7 gün hata oranı (canlı geçiş sonrası)", "<%3", "Haftalık", ">%8 → hypercare uzatma"),
]
for i, row in enumerate(kpi_rows):
    bg = WHITE if i % 2 == 0 else LIGHT_GREY
    table_data_row(slide, row, lefts4, 1.57 + i*0.48, height=0.46, bg=bg)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 17 — QNB DENEYİMİ
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "QNB Deneyimi — Referans Vaka", "Payten'den banka içi sanal POS'a geçiş — 2020, ~2.000 merchant")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# Left — QNB context
add_rect(slide, 0.25, 1.22, 4.5, 5.5, fill_color=WHITE, border_color=DARK_NAVY, border_width=Pt(1.5))
add_rect(slide, 0.25, 1.22, 4.5, 0.45, fill_color=DARK_NAVY)
add_textbox(slide, "🏦  QNB Finansbank, 2020", 0.38, 1.25, 4.2, 0.38, font_size=Pt(11), bold=True, color=WHITE)
qnb_facts = [
    ("Proje Tipi", "Payten → Banka içi\nsanal POS geçişi"),
    ("Portföy", "~2.000 merchant"),
    ("Süre", "~14 ay"),
    ("Sonuç", "Başarıyla tamamlandı"),
]
y = 1.72
for label, value in qnb_facts:
    add_textbox(slide, label + ":", 0.38, y, 1.8, 0.28, font_size=Pt(9.5), bold=True, color=DARK_NAVY)
    add_textbox(slide, value, 2.2, y, 2.4, 0.28, font_size=Pt(9.5), color=DARK_TEXT)
    y += 0.32
    add_rect(slide, 0.38, y-0.02, 4.2, 0.02, fill_color=LIGHT_GREY)

# Right — lessons
add_rect(slide, 4.95, 1.22, 8.1, 2.65, fill_color=WHITE, border_color=ACCENT, border_width=Pt(1.5))
add_rect(slide, 4.95, 1.22, 8.1, 0.42, fill_color=ACCENT)
add_textbox(slide, "Sahadan Çıkan 6 Kritik Ders", 5.08, 1.25, 7.8, 0.36, font_size=Pt(11), bold=True, color=WHITE)
lessons = [
    ("1.", "Toplu e-posta ile migration gerçekleşmiyor — aktif arama + birebir takip şart"),
    ("2.", "Merchant teknik kapasitesi heterojen — tek tip plan başarısız olur"),
    ("3.", "Segmentasyon yapılmadan ekip verimsizleşiyor — önce sınıflandır"),
    ("4.", "Hazır paket kullanıcılar için partner/entegratör kanalı ayrıca yönetilmeli"),
    ("5.", "Merchant karar alma hızı en kritik darboğaz — takip sürekliliği şart"),
    ("6.", "Geçiş sonrası ilk 7 gün destek penceresi başarı için olmazsa olmaz"),
]
y = 1.7
for num, lesson in lessons:
    add_textbox(slide, num, 5.08, y, 0.35, 0.28, font_size=Pt(9.5), bold=True, color=ACCENT)
    add_textbox(slide, lesson, 5.43, y, 7.4, 0.28, font_size=Pt(9.5), color=DARK_TEXT)
    y += 0.3

# X Bank parallel
add_rect(slide, 4.95, 4.05, 8.1, 2.6, fill_color=WHITE, border_color=GREEN, border_width=Pt(1.5))
add_rect(slide, 4.95, 4.05, 8.1, 0.42, fill_color=GREEN)
add_textbox(slide, "X Bank için Bu Deneyimin Değeri", 5.08, 4.08, 7.8, 0.36, font_size=Pt(11), bold=True, color=WHITE)
value_items = [
    "Benzer Payten yapısını daha önce gördük — sürpriz yok",
    "Hazır playbook ve segmentasyon metodolojisi mevcut",
    "Merchant başına gerçek efor tahminlerine sahibiz",
    "Karmaşık segment (ajans bağımlı) için ayrı protokol geliştirildi",
    "İlk gün/ilk hafta destek modelini sahadaki hatalardan öğrendik",
    "X Bank için öngörülü risk yönetimi uygulayabiliriz",
]
y = 4.52
for item in value_items:
    add_textbox(slide, f"✓  {item}", 5.08, y, 7.8, 0.27, font_size=Pt(9.5), color=DARK_TEXT)
    y += 0.28

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 18 — FİNHOUSE KATKISI & FARKLILEŞME
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Finhouse Katkısı & Farklılaşma", "Neden Finhouse — standart outsource'dan farkımız")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

differentiators = [
    ("💳  Ödeme Sistemi\nDomain Uzmanlığı", DARK_NAVY,
     ["Sanal POS, 3D Secure, MPI/ACS, provizyon akışları bilgisi",
      "Türkiye ödeme ekosistemini (Payten, iyzico, BKM) yakından tanıma",
      "BDDK/TCMB mevzuat farkındalığı"]),
    ("🏭  Migration Factory\nMetodolojisi", ACCENT,
     ["Hazır playbook: segment bazlı onboarding protokolü",
      "QNB deneyiminden damıtılmış süreç mimarisi",
      "Wave planı şablonu ve KPI çerçevesi hazır"]),
    ("⚡  Operasyon + Teknik\nTek Çatı", GREEN,
     ["Merchant iletişimi, teknik enablement, PMO ayrı ekipler değil",
      "Banka dilini ve merchant gerçekliğini aynı anda okuma",
      "Koordinasyon maliyeti ve gecikmesi minimize"]),
    ("📊  Ölçülebilir\nDelivery Taahhüdü", ORANGE,
     ["KPI dashboard haftalık yönetim raporuyla sunulur",
      "Discovery fazı çıktısı somut ve test edilebilir",
      "Başarı kriteri net: canlı ilk işlem"]),
]
for i, (title, color, items) in enumerate(differentiators):
    x = 0.25 + (i % 2)*6.55
    y = 1.22 + (i // 2)*2.6
    add_rect(slide, x, y, 6.3, 2.45, fill_color=WHITE, border_color=color, border_width=Pt(1.5))
    add_rect(slide, x, y, 6.3, 0.45, fill_color=color)
    add_textbox(slide, title, x+0.12, y+0.06, 6.05, 0.38, font_size=Pt(11), bold=True, color=WHITE)
    yy = y + 0.52
    for item in items:
        add_textbox(slide, f"▸  {item}", x+0.15, yy, 5.95, 0.28, font_size=Pt(9.5), color=DARK_TEXT)
        yy += 0.3

add_rect(slide, 0.25, 6.62, 12.83, 0.55, fill_color=DARK_NAVY)
add_textbox(slide, "Tek Cümle:  Finhouse, X Bank'a sadece insan kaynağı değil — hazır metodoloji, saha deneyimi ve ölçülebilir sonuç modeli sunar.",
            0.4, 6.72, 12.4, 0.35, font_size=Pt(10), bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 19 — 12 AYLIK YOL HARİTASI (TİMLİNE)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "12 Aylık Yol Haritası", "Ay bazlı migration planı — faz geçişleri ve kilometre taşları")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

# Month header
months = ["A1", "A2", "A3", "A4", "A5", "A6", "A7", "A8", "A9", "A10", "A11", "A12"]
col_w = 1.07
for i, m in enumerate(months):
    x = 0.25 + i*col_w
    bg = DARK_NAVY if i < 2 else (ACCENT if i < 11 else GREEN)
    add_rect(slide, x, 1.22, col_w-0.05, 0.38, fill_color=bg)
    add_textbox(slide, m, x+0.05, 1.25, col_w-0.12, 0.3, font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Activity rows
activities = [
    ("Discovery & Analiz", DARK_NAVY, [1,2]),
    ("Pilot (100-300 merchant)", RGBColor(0x5D, 0x6D, 0x7E), [2,3]),
    ("Wave 1 — Seg. A Hazır Paket", ACCENT, [3,4,5]),
    ("Wave 2 — Seg. B Standart API", RGBColor(0x21, 0x8A, 0xC3), [4,5,6,7]),
    ("Wave 3 — Seg. A+B Kalan", GREEN, [6,7,8]),
    ("Wave 4 — Seg. D KOBİ", ORANGE, [7,8,9]),
    ("Wave 5 — Seg. C Custom/Ajans", RED, [9,10,11]),
    ("Hypercare & Kapanış", RGBColor(0x7D, 0x3C, 0x98), [11,12]),
]
for row_i, (label, color, cols) in enumerate(activities):
    y = 1.72 + row_i*0.46
    add_textbox(slide, label, 0.25, y+0.06, 2.5, 0.32, font_size=Pt(8.5), color=DARK_TEXT)
    for c in range(12):
        x = 0.25 + c*col_w
        if (c+1) in cols:
            add_rect(slide, x+0.05, y+0.05, col_w-0.12, 0.34, fill_color=color)
        else:
            add_rect(slide, x+0.05, y+0.05, col_w-0.12, 0.34, fill_color=LIGHT_GREY, border_color=MID_GREY, border_width=Pt(0.2))

# Milestones
add_textbox(slide, "Kilometre Taşları", 0.25, 5.55, 12.5, 0.3, font_size=Pt(11), bold=True, color=DARK_NAVY)
milestones = [
    ("Ay 2", "Discovery raporu\n+ pilot sonuçları", DARK_NAVY),
    ("Ay 4", "Wave 1 tamamlandı\n~3.000 merchant", ACCENT),
    ("Ay 7", "Portföyün %50'si\ncanlıya geçti", GREEN),
    ("Ay 10", "Portföyün %80'i\nmigre edildi", ORANGE),
    ("Ay 12", "Program kapandı\nfinal rapor", RGBColor(0x7D, 0x3C, 0x98)),
]
for i, (period, label, color) in enumerate(milestones):
    x = 0.25 + i*2.55
    add_rect(slide, x, 5.88, 2.4, 0.8, fill_color=WHITE, border_color=color, border_width=Pt(1.2))
    add_textbox(slide, period, x+0.1, 5.92, 2.2, 0.25, font_size=Pt(9), bold=True, color=color)
    add_textbox(slide, label, x+0.1, 6.17, 2.2, 0.45, font_size=Pt(8.5), color=DARK_TEXT)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 20 — SONRAKI ADIM & CTA
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Sonraki Adım", "Discovery fazına başlamak için 3 adım")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

steps_cta = [
    ("1", "Discovery Toplantısı", "X Bank ile 1 günlük workshop:\nmevcut durum, öncelikler, başarı kriteri",
     "Bu hafta / önümüzdeki hafta", DARK_NAVY),
    ("2", "Ön Veri Seti Paylaşımı", "Merchant listesi, hacim özeti,\nentegrasyon tipi dağılımı",
     "Toplantı sonrası 5 iş günü", ACCENT),
    ("3", "Pilot Kapsamı & Teklif", "100–300 merchant pilot planı,\nnihai SOW ve fiyat netleştirilir",
     "Veri alındıktan 5 iş günü sonra", GREEN),
]
for i, (num, title, desc, timing, color) in enumerate(steps_cta):
    x = 0.5 + i*4.1
    add_rect(slide, x, 1.28, 3.8, 3.2, fill_color=WHITE, border_color=color, border_width=Pt(2))
    add_rect(slide, x, 1.28, 3.8, 0.55, fill_color=color)
    add_textbox(slide, f"ADIM  {num}", x+0.15, 1.32, 3.5, 0.45, font_size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, x+0.15, 1.9, 3.5, 0.38, font_size=Pt(12), bold=True, color=color)
    add_textbox(slide, desc, x+0.15, 2.35, 3.5, 0.7, font_size=Pt(10), color=DARK_TEXT)
    add_rect(slide, x+0.15, 3.12, 3.5, 0.04, fill_color=color)
    add_textbox(slide, f"⏱  {timing}", x+0.15, 3.2, 3.5, 0.32, font_size=Pt(9.5), color=color, italic=True)

# Contact & CTA
add_rect(slide, 0.25, 4.65, 12.83, 1.65, fill_color=DARK_NAVY)
add_textbox(slide, "İletişim & Bir Sonraki Toplantı",
            0.4, 4.72, 12.3, 0.35, font_size=Pt(13), bold=True, color=WHITE)
add_textbox(slide, "Ayhan Ağırgöl  |  Finhouse Danışmanlık",
            0.4, 5.12, 5.5, 0.32, font_size=Pt(11), color=ACCENT)
add_textbox(slide, "ayhan.agirgol@finhouse.com.tr",
            0.4, 5.44, 5.5, 0.28, font_size=Pt(10), color=RGBColor(0xCC, 0xDD, 0xEE))
add_textbox(slide, "finhouse.ai",
            0.4, 5.72, 3, 0.28, font_size=Pt(10), color=ACCENT)
add_textbox(slide, "\"18.000 merchant'ı 12 ayda taşımak — sahadan gelen metodoloji, ölçülebilir sonuç.\"",
            6.5, 5.05, 6.3, 0.6, font_size=Pt(11), italic=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 21 — EKLER / REFERANS VERİLER
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_header(slide, "Ekler — Referans Veri & Kaynaklar", "Araştırma metodolojisi ve veri güvenilirliği notu")
add_rect(slide, 0, 1.09, 13.33, 6.06, fill_color=LIGHT_GREY)

add_textbox(slide, "Doğrulanan Kaynaklar", 0.25, 1.2, 6, 0.32, font_size=Pt(12), bold=True, color=DARK_NAVY)
sources = [
    "Stripe Connect Migration Docs — docs.stripe.com/connect/migrate-to-controller-properties",
    "Stripe Pricing — stripe.com/pricing (kamuya açık, doğrulandı)",
    "Adyen Pricing — adyen.com/pricing (Interchange++ model, doğrulandı)",
    "BKM 2024 Yıllık İstatistikleri — bkm.com.tr/istatistikler",
    "TCMB Ödeme İstatistikleri — tcmb.gov.tr",
    "BDDK Ödeme Hizmetleri Lisans Veritabanı — bddk.org.tr",
    "iyzico Ürün Sayfası — iyzico.com (Türkiye yerel benchmark)",
]
y = 1.58
for s in sources:
    add_textbox(slide, f"✓  {s}", 0.25, y, 6.4, 0.27, font_size=Pt(9), color=DARK_TEXT)
    y += 0.28

add_textbox(slide, "Sektör Verisi (Yaklaşık / Kamuya Açık Değil)", 6.8, 1.2, 6.2, 0.32, font_size=Pt(12), bold=True, color=ORANGE)
approx_sources = [
    "Payten/Asseco SEE fiyat yapısı — sektör müzakere bilgisi",
    "Worldpay/FIS, PayPal/Braintree migration süreleri — sektör raporu",
    "Türkiye yerel banka Payten→kendi gateway geçişi — sektör içi bilgi",
    "McKinsey Global Payments 2023 — erişim engellendi, genel bilinirlik",
    "MDR bantları (iyzico, Paynet, banka) — münferit sözleşme gözlemi",
    "Merchant churn %15–30 riski — sektör tahmini",
]
y = 1.58
for s in approx_sources:
    add_textbox(slide, f"⚠  {s}", 6.8, y, 6.2, 0.27, font_size=Pt(9), color=DARK_TEXT)
    y += 0.28

# Glossary
add_textbox(slide, "Terimler Sözlüğü", 0.25, 3.62, 12.5, 0.3, font_size=Pt(11), bold=True, color=DARK_NAVY)
terms = [
    ("MDR", "Merchant Discount Rate — işlem komisyon oranı"),
    ("KYC / MASAK", "Müşterini Tanı / suç gelirlerinin aklanmasını önleme kontrolü"),
    ("3D Secure", "Kartlı ödeme güvenlik protokolü (Visa Secure / MC ID Check)"),
    ("PCI-DSS SAQ-A", "Ödeme güvenlik standardı — hosted page kullanıcılar için"),
    ("Wave Migration", "Merchant'ların segment ve öncelik sırasıyla dalgalar halinde taşınması"),
    ("Never Dark", "Geçiş sürecinde hiçbir merchant'ın işlem yapamamaz duruma düşmemesi"),
    ("IC++ (Interchange++)", "Kart ağı maliyeti + banka marjı + işlemci ücreti şeffaf ayrıştırma modeli"),
    ("Hypercare", "Canlı geçiş sonrası ilk 30 gün yoğun teknik destek penceresi"),
]
col1 = terms[:4]
col2 = terms[4:]
y = 4.0
for (term, desc) in col1:
    add_textbox(slide, f"{term}:", 0.25, y, 1.7, 0.27, font_size=Pt(9), bold=True, color=DARK_NAVY)
    add_textbox(slide, desc, 1.95, y, 4.5, 0.27, font_size=Pt(9), color=DARK_TEXT)
    y += 0.28
y = 4.0
for (term, desc) in col2:
    add_textbox(slide, f"{term}:", 6.8, y, 1.7, 0.27, font_size=Pt(9), bold=True, color=DARK_NAVY)
    add_textbox(slide, desc, 8.5, y, 4.6, 0.27, font_size=Pt(9), color=DARK_TEXT)
    y += 0.28

# ════════════════════════════════════════════════════════════════════════════
#  SLIDE 22 — BACK COVER
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_NAVY)
add_rect(slide, 0, 5.5, 13.33, 0.06, fill_color=ACCENT)
add_rect(slide, 0, 5.56, 13.33, 1.94, fill_color=RGBColor(0x12, 0x1E, 0x35))

add_textbox(slide, "finhouse", 0.7, 1.8, 5, 0.9, font_size=Pt(48), bold=True, color=WHITE)
add_textbox(slide, ".ai", 4.3, 1.8, 2, 0.9, font_size=Pt(48), bold=True, color=ACCENT)
add_textbox(slide, "Fintech & AI Çözümleri  |  Danışmanlık  |  Operasyon",
            0.7, 2.8, 7, 0.45, font_size=Pt(14), color=RGBColor(0xAA, 0xCC, 0xEE))
add_rect(slide, 0.7, 3.45, 5, 0.04, fill_color=ACCENT)
add_textbox(slide, "ayhan.agirgol@finhouse.com.tr", 0.7, 3.65, 5, 0.35, font_size=Pt(11), color=WHITE)
add_textbox(slide, "www.finhouse.ai", 0.7, 4.0, 4, 0.35, font_size=Pt(11), color=ACCENT)

add_textbox(slide, "Bu doküman X Bank yetkilileri için hazırlanmıştır.\nGizli — Üçüncü şahıslarla paylaşılamaz.",
            0.7, 5.7, 8, 0.65, font_size=Pt(9), italic=True, color=RGBColor(0x88, 0xAA, 0xCC))

# Slide count
add_textbox(slide, f"Toplam: 22 slayt  |  Finhouse © 2026", 10.5, 7.18, 2.5, 0.25,
            font_size=Pt(8), color=MID_GREY, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════════════════
#  SAVE
# ════════════════════════════════════════════════════════════════════════════
output_path = "/Users/baykus/.openclaw/workspace/xbank_migration/xbank_detailed_client_deck.pptx"
prs.save(output_path)
print(f"✅  Presentation saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
